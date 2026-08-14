#!/usr/bin/env python3
"""Build an MSU-branded .pptx from a markdown deck spec.

Usage:
    python build_deck.py <spec.md> -o <out.pptx> [--template T] [--tokens dist/msu.css]
    python build_deck.py <spec.md> --validate       # parse + report warnings, write nothing

Design constraints, and why:

1. NO BRAND COLOR LITERALS. The MSU brand kit permits a hand-written brand
   color literal in exactly two of its own files. This is not one of them.
   By default this script writes no color values at all: it references the
   template's theme slots (accent1/accent2), so whatever the MSU template
   defines is what appears. With --tokens it parses published values out of
   the kit's GENERATED dist/msu.css instead. Either way the values are
   referenced, never retyped here.

2. The MSU logo is a blipFill on the *layout's* placeholder, not a slide
   shape. python-pptx clones layout placeholders onto each new slide without
   that fill, so an unfilled clone hides the logo. We therefore delete every
   cloned placeholder we did not populate, which is exactly what MSU's own
   sample slides do.

3. Slides are a screen target, so web color values apply, never print.
"""

from __future__ import annotations

import argparse
import re
import sys
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

# Layout names in BlankMSUPowerPoint_*.pptx. Looked up by name, with an index
# fallback, so a template revision that reorders layouts does not break us.
LAYOUTS = {
    "title": ("Title Slide", 0),
    "content": ("Title and Content", 1),
    "section": ("Section Header", 2),
    "two": ("Two Content", 3),
    "titleonly": ("Title Only", 4),
}

# Content region of the MSU 16:9 layouts, read off the template's own
# placeholders rather than chosen by us.
BODY_LEFT = Emu(838200)
BODY_WIDTH = Emu(10515600)
BODY_TOP = Emu(1986243)
BODY_BOTTOM = Emu(6172200)

MONO = "Consolas"  # local choice for Excel formulas; not an MSU standard

# Bullet-count thresholds -> body point size. A teaching slide past ~12 lines
# is a handout, not a slide, so we warn rather than silently shrink forever.
DENSITY_STEPS = ((5, None), (8, 18), (11, 16), (14, 14))


# --------------------------------------------------------------------------
# spec parsing
# --------------------------------------------------------------------------

class Slide:
    def __init__(self, kind, title=""):
        self.kind = kind          # title | section | content
        self.title = title
        self.subheading = ""
        self.blocks = []          # ("bullet", text, level) | ("para", text, 0) | ("quote", text, 0)
        self.table = None         # list[list[str]], first row is the header
        self.left = []            # two-column mode: blocks for each side
        self.right = []
        self.notes = ""

    @property
    def two_col(self):
        return bool(self.left or self.right)


def _parse_meta(lines, i):
    """Read a --- fenced key: value block if one opens the file."""
    meta = {}
    if i >= len(lines) or lines[i].strip() != "---":
        return meta, i
    i += 1
    while i < len(lines) and lines[i].strip() != "---":
        if ":" in lines[i]:
            k, v = lines[i].split(":", 1)
            meta[k.strip().lower()] = v.strip()
        i += 1
    return meta, i + 1


def parse_spec(text):
    lines = text.replace("\r\n", "\n").split("\n")
    meta, i = _parse_meta(lines, 0)

    slides = []
    cur = None
    col = None          # None | "left" | "right"
    notes_mode = False
    table_rows = None

    def flush_table():
        nonlocal table_rows
        if table_rows and cur is not None:
            # drop the |---|---| separator row
            rows = [r for r in table_rows if not re.fullmatch(r"[\s|:-]+", "|".join(r))]
            if rows:
                cur.table = rows
        table_rows = None

    def target():
        if col == "left":
            return cur.left
        if col == "right":
            return cur.right
        return cur.blocks

    for raw in lines[i:]:
        line = raw.rstrip()
        stripped = line.strip()

        if stripped.startswith("|") and cur is not None:
            cells = [c.strip() for c in stripped.strip("|").split("|")]
            if table_rows is None:
                table_rows = []
            table_rows.append(cells)
            continue
        flush_table()

        if not stripped:
            if notes_mode and cur is not None and cur.notes:
                cur.notes += "\n"
            continue

        m = re.match(r"^(#{1,3})\s+(.*)$", stripped)
        if m:
            hashes, txt = m.group(1), m.group(2).strip()
            if hashes == "###":
                if cur is not None:
                    cur.subheading = txt
                continue
            cur = Slide("section" if hashes == "#" else "content", txt)
            slides.append(cur)
            col = None
            notes_mode = False
            continue

        if cur is None:
            continue

        if stripped.startswith(":::"):
            side = stripped[3:].strip().lower()
            col = side if side in ("left", "right") else None
            notes_mode = False
            continue

        m = re.match(r"^(NOTES:|\?\?\?)\s*(.*)$", stripped)
        if m:
            notes_mode = True
            cur.notes = m.group(2)
            continue

        if notes_mode:
            cur.notes += ("\n" if cur.notes else "") + stripped
            continue

        m = re.match(r"^(\s*)[-*+]\s+(.*)$", line)
        if m:
            level = min(len(m.group(1)) // 2, 4)
            target().append(("bullet", m.group(2).strip(), level))
            continue

        if stripped.startswith(">"):
            target().append(("quote", stripped.lstrip("> ").strip(), 0))
            continue

        target().append(("para", stripped, 0))

    flush_table()
    return meta, slides


# --------------------------------------------------------------------------
# inline formatting: **bold**, *italic*, `code`
# --------------------------------------------------------------------------

_INLINE = re.compile(r"(\*\*.+?\*\*|(?<!\*)\*[^*]+?\*(?!\*)|`[^`]+?`)")


def write_runs(paragraph, text, mono_name=MONO):
    for piece in _INLINE.split(text):
        if not piece:
            continue
        run = paragraph.add_run()
        if piece.startswith("**") and piece.endswith("**"):
            run.text = piece[2:-2]
            run.font.bold = True
        elif piece.startswith("`") and piece.endswith("`"):
            run.text = piece[1:-1]
            run.font.name = mono_name
        elif piece.startswith("*") and piece.endswith("*"):
            run.text = piece[1:-1]
            run.font.italic = True
        else:
            run.text = piece


# --------------------------------------------------------------------------
# color: theme references by default, kit tokens on request
# --------------------------------------------------------------------------

class Palette:
    """Either theme-slot references (no literals) or values parsed from the
    brand kit's generated dist/msu.css. Never values typed into this file."""

    def __init__(self, css_path=None):
        self.blue = self.gold = self.white = None
        self.source = "template theme (accent1 / accent2 / lt1)"
        if css_path:
            css = Path(css_path).read_text(encoding="utf-8")

            def var(name):
                m = re.search(rf"--{name}:\s*(#[0-9a-fA-F]{{6}})\s*;", css)
                if not m:
                    raise SystemExit(f"--tokens: {name} not found in {css_path}")
                return RGBColor.from_string(m.group(1)[1:])

            # Slides are a screen target: web values, never print.
            self.blue = var("msu-web-blue")
            self.gold = var("msu-web-gold")
            self.white = var("kit-neutral-white")
            self.source = f"{css_path} (msu-web-*)"

    def apply_fill(self, fill, which):
        fill.solid()
        if self.blue is None:
            fill.fore_color.theme_color = (
                MSO_THEME_COLOR.ACCENT_1 if which == "blue" else MSO_THEME_COLOR.ACCENT_2
            )
        else:
            fill.fore_color.rgb = self.blue if which == "blue" else self.gold

    def apply_text_white(self, font):
        if self.white is None:
            font.color.theme_color = MSO_THEME_COLOR.BACKGROUND_1
        else:
            font.color.rgb = self.white


# --------------------------------------------------------------------------
# pptx helpers
# --------------------------------------------------------------------------

def strip_slides(prs):
    """Remove the template's own sample slides."""
    id_list = prs.slides._sldIdLst
    for sld_id in list(id_list):
        prs.part.drop_rel(sld_id.get(qn("r:id")))
        id_list.remove(sld_id)


def layout(prs, key):
    name, idx = LAYOUTS[key]
    for lay in prs.slide_layouts:
        if lay.name == name:
            return lay
    return prs.slide_layouts[idx]


def send_to_back(shape):
    tree = shape._element.getparent()
    tree.remove(shape._element)
    # index 2 clears <p:nvGrpSpPr> and <p:grpSpPr>
    tree.insert(2, shape._element)


def drop_unused_placeholders(slide):
    """Delete cloned placeholders we left empty.

    Required, not cosmetic: layout placeholders such as 'MSU Logo' carry the
    logo as a picture fill that cloning does not copy, so an empty clone
    covers the layout's logo with nothing.
    """
    for ph in list(slide.placeholders):
        if ph.has_text_frame and ph.text_frame.text.strip():
            continue
        if ph.shape_type == 13:  # a picture we inserted
            continue
        ph._element.getparent().remove(ph._element)


def body_size(blocks):
    n = sum(1 for b in blocks if b[0] in ("bullet", "para", "quote"))
    for limit, size in DENSITY_STEPS:
        if n <= limit:
            return size
    return 12


def fill_text_frame(tf, blocks, size, palette):
    tf.word_wrap = True
    first = True
    for kind, text, level in blocks:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.level = level
        if kind == "quote":
            para.alignment = PP_ALIGN.LEFT
        write_runs(para, text)
        if kind == "quote":
            for run in para.runs:
                run.font.italic = True
        if size:
            for run in para.runs:
                run.font.size = Pt(size if level == 0 else max(size - 2, 10))


def add_table(slide, rows, top, height, palette, warn):
    n_rows, n_cols = len(rows), max(len(r) for r in rows)
    rows = [r + [""] * (n_cols - len(r)) for r in rows]
    shape = slide.shapes.add_table(n_rows, n_cols, BODY_LEFT, top, BODY_WIDTH, height)
    table = shape.table
    size = 14 if n_cols <= 3 else (12 if n_cols <= 5 else 11)
    if n_cols > 6:
        warn(f"table has {n_cols} columns; over 6 rarely reads from the back of JABS 215")

    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = ""
            para = cell.text_frame.paragraphs[0]
            write_runs(para, val)
            for run in para.runs:
                run.font.size = Pt(size)
                if r == 0:
                    run.font.bold = True
                    palette.apply_text_white(run.font)
            if r == 0:
                palette.apply_fill(cell.fill, "blue")
    return shape


# --------------------------------------------------------------------------
# slide builders
# --------------------------------------------------------------------------

def build_title_slide(prs, meta, branding, palette):
    slide = prs.slides.add_slide(layout(prs, "title"))
    if slide.shapes.title is not None:
        slide.shapes.title.text_frame.text = meta.get("title", "")
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            ph.text_frame.text = meta.get("subtitle", "")

    # Re-create the template's own title-slide furniture: photo, shade, logo.
    for name, blob, left, top, width, height in branding:
        pic = slide.shapes.add_picture(BytesIO(blob), left, top, width, height)
        send_to_back(pic)
    if branding:  # text sits on a shaded photo
        for ph in slide.placeholders:
            if ph.has_text_frame:
                for para in ph.text_frame.paragraphs:
                    for run in para.runs:
                        palette.apply_text_white(run.font)
    drop_unused_placeholders(slide)
    return slide


def build_slide(prs, spec, palette, warn):
    if spec.kind == "section":
        slide = prs.slides.add_slide(layout(prs, "section"))
        if slide.shapes.title is not None:
            slide.shapes.title.text_frame.text = spec.title
        body = [b for b in spec.blocks]
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1 and body:
                fill_text_frame(ph.text_frame, body, None, palette)
        drop_unused_placeholders(slide)
        _attach_notes(slide, spec)
        return slide

    has_body = bool(spec.blocks)
    if spec.two_col:
        slide = prs.slides.add_slide(layout(prs, "two"))
    elif spec.table is not None and not has_body:
        slide = prs.slides.add_slide(layout(prs, "titleonly"))
    elif spec.table is not None and has_body:
        slide = prs.slides.add_slide(layout(prs, "titleonly"))
    else:
        slide = prs.slides.add_slide(layout(prs, "content"))

    if slide.shapes.title is not None:
        slide.shapes.title.text_frame.text = spec.title
    if spec.subheading:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 20:
                ph.text_frame.text = spec.subheading
                break
        else:
            warn(f"'{spec.title}': layout has no Subheading placeholder; subheading dropped")

    top = BODY_TOP if not spec.subheading else Emu(int(BODY_TOP) + 300000)

    if spec.two_col:
        size = min((s for s in (body_size(spec.left), body_size(spec.right)) if s),
                   default=None)
        targets = [ph for ph in slide.placeholders
                   if ph.placeholder_format.idx in (1, 2)]
        targets.sort(key=lambda p: p.left)
        for ph, blocks in zip(targets, (spec.left, spec.right)):
            if blocks:
                fill_text_frame(ph.text_frame, blocks, size, palette)
    elif spec.table is not None and has_body:
        # bullets above, table below
        split = Emu(int(top) + 1200000)
        box = slide.shapes.add_textbox(BODY_LEFT, top, BODY_WIDTH, Emu(1100000))
        fill_text_frame(box.text_frame, spec.blocks, body_size(spec.blocks) or 18, palette)
        add_table(slide, spec.table, split, Emu(int(BODY_BOTTOM) - int(split)), palette, warn)
    elif spec.table is not None:
        add_table(slide, spec.table, top, Emu(int(BODY_BOTTOM) - int(top)), palette, warn)
    else:
        target = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                target = ph
                break
        if target is None:
            target = slide.shapes.add_textbox(
                BODY_LEFT, top, BODY_WIDTH, Emu(int(BODY_BOTTOM) - int(top)))
        fill_text_frame(target.text_frame, spec.blocks, body_size(spec.blocks), palette)

    drop_unused_placeholders(slide)
    _attach_notes(slide, spec)
    return slide


def _attach_notes(slide, spec):
    if spec.notes.strip():
        slide.notes_slide.notes_text_frame.text = spec.notes.strip()


# --------------------------------------------------------------------------
# quality warnings: a deck the room can actually read
# --------------------------------------------------------------------------

def lint(slides, warn):
    for s in slides:
        n = sum(1 for b in s.blocks if b[0] in ("bullet", "para", "quote"))
        if n > 14:
            warn(f"'{s.title}': {n} lines on one slide — split it")
        if len(s.title) > 70:
            warn(f"title runs {len(s.title)} chars and will wrap: '{s.title[:50]}...'")
        for kind, text, _ in s.blocks:
            if kind == "bullet" and len(text) > 130:
                warn(f"'{s.title}': a bullet runs {len(text)} chars — that is a sentence "
                     f"for the notes field, not a bullet")
    if not any(s.notes.strip() for s in slides):
        warn("no slide carries NOTES: — the instructor moves live in the lesson plan, "
             "but a deck with no notes is hard to teach from a year later")


# --------------------------------------------------------------------------

def main(argv=None):
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("spec", help="markdown deck spec")
    ap.add_argument("-o", "--out", help="output .pptx")
    ap.add_argument("--template", default="Branding/BlankMSUPowerPoint_2026.04.15.pptx")
    ap.add_argument("--tokens", metavar="MSU_CSS",
                    help="path to the brand kit's generated dist/msu.css; uses published "
                         "msu-web-* values instead of the template's theme slots")
    ap.add_argument("--validate", action="store_true", help="parse and lint only")
    args = ap.parse_args(argv)

    spec_path = Path(args.spec)
    meta, slides = parse_spec(spec_path.read_text(encoding="utf-8"))

    problems = []
    warn = problems.append
    if not meta.get("title"):
        warn("spec has no 'title:' in its frontmatter")
    lint(slides, warn)

    if args.validate:
        print(f"{spec_path}: {len(slides)} slides parsed")
        for p in problems:
            print(f"  warning: {p}")
        return 0 if not problems else 0

    if not args.out:
        ap.error("-o/--out is required unless --validate")

    template = Path(args.template)
    if not template.exists():
        raise SystemExit(f"template not found: {template}")

    palette = Palette(args.tokens)
    prs = Presentation(str(template))

    # Grab the template's title-slide furniture before dropping its slides.
    branding = []
    title_name = LAYOUTS["title"][0]
    for s in prs.slides:
        if s.slide_layout.name == title_name:
            for sh in s.shapes:
                if sh.shape_type == 13:
                    branding.append((sh.name, sh.image.blob,
                                     sh.left, sh.top, sh.width, sh.height))
            break

    strip_slides(prs)
    build_title_slide(prs, meta, branding, palette)
    for spec in slides:
        build_slide(prs, spec, palette, warn)

    out = Path(args.out)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))

    print(f"wrote {out}  ({len(prs.slides._sldIdLst)} slides)")
    print(f"  template: {template}")
    print(f"  colors:   {palette.source}")
    for p in problems:
        print(f"  warning: {p}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
